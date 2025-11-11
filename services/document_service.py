"""
Document Processing Service
Extracts text content from PDF and PowerPoint files
"""

import logging
from typing import Optional
import io

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def extract_text_from_pdf(file_content: bytes) -> str:
    """
    Extract text from PDF file
    
    Args:
        file_content (bytes): PDF file content
        
    Returns:
        str: Extracted text content
    """
    try:
        from PyPDF2 import PdfReader
        
        pdf_file = io.BytesIO(file_content)
        reader = PdfReader(pdf_file)
        
        text_content = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_content.append(text)
        
        extracted_text = '\n\n'.join(text_content)
        logger.info(f"Extracted {len(extracted_text)} characters from PDF")
        
        return extracted_text
        
    except Exception as e:
        logger.error(f"PDF extraction error: {e}")
        raise Exception(f"Failed to extract PDF content: {str(e)}")


def extract_text_from_pptx(file_content: bytes) -> str:
    """
    Extract text from PowerPoint file
    
    Args:
        file_content (bytes): PowerPoint file content
        
    Returns:
        str: Extracted text content
    """
    try:
        from pptx import Presentation
        
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)
        
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            
            # Extract text from tables
            if hasattr(slide, 'shapes'):
                for shape in slide.shapes:
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text)
                            if row_text:
                                slide_text.append(' | '.join(row_text))
            
            if len(slide_text) > 1:  # Has content beyond the header
                text_content.append('\n'.join(slide_text))
        
        extracted_text = '\n\n'.join(text_content)
        logger.info(f"Extracted {len(extracted_text)} characters from PPTX ({len(prs.slides)} slides)")
        
        return extracted_text
        
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        raise Exception(f"Failed to extract PowerPoint content: {str(e)}")


def extract_document_content(filename: str, file_content: bytes) -> str:
    """
    Extract text content from document based on file extension
    
    Args:
        filename (str): Name of the file
        file_content (bytes): File content
        
    Returns:
        str: Extracted text content
    """
    filename_lower = filename.lower()
    
    if filename_lower.endswith('.pdf'):
        return extract_text_from_pdf(file_content)
    elif filename_lower.endswith(('.ppt', '.pptx')):
        return extract_text_from_pptx(file_content)
    else:
        raise ValueError(f"Unsupported file format: {filename}. Only PDF and PowerPoint files are supported.")


def summarize_content(text: str, max_length: int = 3000) -> str:
    """
    Intelligently summarize or truncate content to fit within AI token limits
    Prioritizes keeping the most important content from beginning, middle, and end
    
    Args:
        text (str): Full text content
        max_length (int): Maximum character length
        
    Returns:
        str: Summarized or truncated text
    """
    if len(text) <= max_length:
        return text
    
    # For very long content, take samples from beginning, middle, and end
    # to ensure comprehensive coverage
    if len(text) > max_length * 3:
        # Take 40% from beginning, 30% from middle, 30% from end
        start_length = int(max_length * 0.4)
        middle_length = int(max_length * 0.3)
        end_length = int(max_length * 0.3)
        
        start_text = text[:start_length]
        middle_start = len(text) // 2 - middle_length // 2
        middle_text = text[middle_start:middle_start + middle_length]
        end_text = text[-end_length:]
        
        # Clean up at sentence boundaries
        start_text = _cut_at_boundary(start_text, from_end=True)
        middle_text = _cut_at_boundary(middle_text, from_start=True)
        middle_text = _cut_at_boundary(middle_text, from_end=True)
        end_text = _cut_at_boundary(end_text, from_start=True)
        
        return f"{start_text}\n\n[... content from middle section ...]\n\n{middle_text}\n\n[... content from end section ...]\n\n{end_text}"
    
    # For moderately long content, just truncate at a good boundary
    truncated = text[:max_length]
    truncated = _cut_at_boundary(truncated, from_end=True)
    
    return truncated + "\n\n... (content continues)"


def _cut_at_boundary(text: str, from_start: bool = False, from_end: bool = False) -> str:
    """
    Cut text at a natural boundary (paragraph, sentence, or space)
    
    Args:
        text (str): Text to cut
        from_start (bool): If True, cut from the start to first boundary
        from_end (bool): If True, cut from the end to last boundary
        
    Returns:
        str: Text cut at boundary
    """
    if not text:
        return text
    
    boundaries = ['\n\n', '\n', '. ', '。', '! ', '！', '? ', '？', ' ']
    
    if from_end:
        # Find last boundary and cut there
        for delimiter in boundaries:
            last_index = text.rfind(delimiter)
            if last_index > len(text) * 0.85:  # Within last 15%
                return text[:last_index + len(delimiter)]
        return text
    
    if from_start:
        # Find first boundary and cut there
        for delimiter in boundaries:
            first_index = text.find(delimiter)
            if first_index > 0 and first_index < len(text) * 0.15:  # Within first 15%
                return text[first_index + len(delimiter):]
        return text
    
    return text
